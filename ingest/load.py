"""
U1: Document loading module.
load_pdf(path: str) -> list[tuple[str, int]]
"""
import sys
from pathlib import Path
from typing import List, Tuple, Optional

if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass


def load_pdf(path: str) -> List[Tuple[str, int]]:
    """
    Load a PDF and return a list of (text, page_number) pairs.
    1-indexed. Every page produces an entry, even blank pages ('', page_num).
    """
    import fitz  # PyMuPDF

    file_path = Path(path)
    if not file_path.exists():
        raise FileNotFoundError(f"PDF file not found: {path}")

    doc = fitz.open(str(file_path))
    pages: List[Tuple[str, int]] = []

    for page_idx in range(len(doc)):
        page_num = page_idx + 1
        page = doc[page_idx]
        text = page.get_text("text").strip()
        pages.append((text, page_num))

    doc.close()
    return pages


def load_docx(path: str) -> List[Tuple[str, Optional[int]]]:
    """Load a DOCX file and return (text, None) pairs."""
    import docx

    file_path = Path(path)
    if not file_path.exists():
        raise FileNotFoundError(f"DOCX file not found: {path}")

    doc = docx.Document(str(file_path))
    full_paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    full_text = "\n\n".join(full_paragraphs)
    return [(full_text, None)]


def load_pptx(path: str) -> List[Tuple[str, Optional[int]]]:
    """Load a PPTX file and return (slide_text, slide_number) pairs."""
    import pptx

    file_path = Path(path)
    if not file_path.exists():
        raise FileNotFoundError(f"PPTX file not found: {path}")

    prs = pptx.Presentation(str(file_path))
    pages: List[Tuple[str, Optional[int]]] = []

    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        texts.append(t)
        pages.append(("\n".join(texts), idx))

    return pages


def load_txt(path: str) -> List[Tuple[str, Optional[int]]]:
    """Load a plain TXT file. Returns [(content, None)]."""
    file_path = Path(path)
    if not file_path.exists():
        raise FileNotFoundError(f"TXT file not found: {path}")

    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read().strip()
    return [(content, None)]


def load_document(path: str) -> List[Tuple[str, Optional[int]]]:
    """Generic document loader routing to the appropriate format handler."""
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return [(t, p) for t, p in load_pdf(path)]
    elif ext == ".docx":
        return load_docx(path)
    elif ext == ".pptx":
        return load_pptx(path)
    else:
        return load_txt(path)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python -m ingest.load <path_to_pdf>")
        sys.exit(1)

    target_path = sys.argv[1]
    results = load_pdf(target_path)
    print(f"Loaded {len(results)} page(s) from {target_path}:\n")
    for text, page_num in results:
        preview = text[:80].replace("\n", " ").strip()
        print(f"page {page_num}: {preview if preview else '<empty page>'}")
